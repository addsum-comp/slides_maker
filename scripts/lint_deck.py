#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""lint_deck.py — deterministic LAYOUT lint for a built .pptx: catches the overlaps a visual
critic can miss (a callout tucked a few px under a panel/image, a block hanging off the slide,
content colliding with the footer).

    python scripts/lint_deck.py deck.pptx

Run it right after building, BEFORE/with the critic loop — it's a cheap, deterministic safety net
for the "no block/text/image overlap" rule, not a replacement for the visual critic (which judges
crop, balance, legibility, fidelity). Exits non-zero if it finds anything.

Checks (tuned for low false-positives):
  1. OVERFLOW   — a shape extends beyond the slide edge.
  2. OVERLAP    — two SOLID shapes (filled auto-shape / picture / table / freeform) partially
                  overlap with NEITHER contained in the other. Intentional layering — text on a
                  card, an image inside its frame, a header band on a card, a full-bleed
                  background — is *containment* and is NOT flagged; only "two separate blocks
                  collide" partial overlaps are.
  3. FOOTER     — a solid block overlaps a footer / page-number text box at the bottom.
"""
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU = 914400.0
TOL = 0.05        # inches — ignore hairline/touching overlaps
CONTAIN = 0.90    # A is "inside" B when >=90% of A's area lies within B
SOLID = {MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.TABLE, MSO_SHAPE_TYPE.FREEFORM}


def _boxes(slide, sw, sh):
    out = []
    for s in slide.shapes:
        try:
            l, t, w, h = s.left / EMU, s.top / EMU, s.width / EMU, s.height / EMU
        except (TypeError, AttributeError):
            continue
        if not w or not h or w <= 0 or h <= 0:
            continue
        txt = s.text_frame.text.strip().replace("\n", " ")[:26] if s.has_text_frame else ""
        out.append({"l": l, "t": t, "w": w, "h": h, "r": l + w, "b": t + h,
                    "st": str(s.shape_type).split()[0], "txt": txt,
                    "solid": s.shape_type in SOLID,
                    "text": bool(s.has_text_frame and txt),
                    "bg": (w * h) >= 0.95 * (sw * sh)})
    return out


def _inter(a, b):
    return (max(0.0, min(a["r"], b["r"]) - max(a["l"], b["l"])),
            max(0.0, min(a["b"], b["b"]) - max(a["t"], b["t"])))


def _frac_inside(a, b):
    ix, iy = _inter(a, b)
    return (ix * iy) / (a["w"] * a["h"] + 1e-9)


def lint(path):
    prs = Presentation(path)
    sw, sh = prs.slide_width / EMU, prs.slide_height / EMU
    total = 0
    for si, slide in enumerate(prs.slides):
        bx = _boxes(slide, sw, sh)
        finds = []
        # 1) overflow
        for s in bx:
            if s["bg"]:
                continue
            ov = []
            if s["l"] < -TOL: ov.append("left")
            if s["t"] < -TOL: ov.append("top")
            if s["r"] > sw + TOL: ov.append(f"right+{round(s['r']-sw,2)}")
            if s["b"] > sh + TOL: ov.append(f"bottom+{round(s['b']-sh,2)}")
            if ov:
                finds.append(f"OVERFLOW [{','.join(ov)}] {s['st']} '{s['txt']}'")
        # 2) solid vs solid partial overlap (neither contained)
        sol = [s for s in bx if s["solid"] and not s["bg"]]
        for i in range(len(sol)):
            for j in range(i + 1, len(sol)):
                a, b = sol[i], sol[j]
                ix, iy = _inter(a, b)
                if ix > TOL and iy > TOL and _frac_inside(a, b) < CONTAIN and _frac_inside(b, a) < CONTAIN:
                    finds.append(f"OVERLAP {round(ix,2)}x{round(iy,2)}in  {a['st']}'{a['txt']}' x {b['st']}'{b['txt']}'")
        # 3) footer collision: a solid block over a bottom text label
        footers = [s for s in bx if s["text"] and s["t"] > sh - 0.6 and not s["bg"]]
        for f in footers:
            for s in sol:
                ix, iy = _inter(s, f)
                if ix > TOL and iy > TOL and _frac_inside(f, s) < CONTAIN:
                    finds.append(f"FOOTER collision  {s['st']}'{s['txt']}' over footer '{f['txt']}'")
        for m in finds:
            print(f"  slide {si+1}: {m}")
        total += len(finds)
    print(f"\n{path}: {total} layout finding(s)" + ("" if total else "  ✓ clean"))
    return total


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("usage: python lint_deck.py <deck.pptx>"); sys.exit(2)
    sys.exit(1 if lint(sys.argv[1]) > 0 else 0)
